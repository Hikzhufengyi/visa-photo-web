from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/seafon/Downloads/北京师范大学_调研汇报模板.pptx")
LOGO = Path("/Users/seafon/Downloads/275dc16574c72ad273e6b03311fae1c7.PNG")

BLUE = RGBColor(0, 91, 172)
DARK = RGBColor(34, 43, 54)
MID = RGBColor(100, 110, 122)
PALE = RGBColor(248, 250, 252)
SOFT = RGBColor(234, 241, 249)
WHITE = RGBColor(255, 255, 255)
LINE = RGBColor(0, 91, 172)


def text(slide, x, y, w, h, value, size=18, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, values, size=16, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, value in enumerate(values):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = value
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
    return box


def logo(slide, size=0.68):
    slide.shapes.add_picture(str(LOGO), Inches(12.05), Inches(0.27), width=Inches(size))


def rule(slide, y):
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.72), Inches(y), Inches(11.28), Inches(0.006)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()
    return line


def header(slide, title, page=None):
    logo(slide)
    text(slide, 0.72, 0.48, 8.0, 0.38, title, 21, BLUE, True)
    rule(slide, 1.05)
    rule(slide, 6.86)
    if page:
        text(slide, 12.15, 7.05, 0.45, 0.2, str(page), 8.5, MID, align=PP_ALIGN.RIGHT)


def pale_box(slide, x, y, w, h):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PALE
    box.line.fill.background()
    return box


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_picture(str(LOGO), Inches(11.62), Inches(0.36), width=Inches(0.86))
    rule(slide, 1.08)
    rule(slide, 6.86)
    text(slide, 0.78, 2.18, 9.6, 0.72, "调研汇报题目", 35, BLUE, True)
    text(slide, 0.8, 3.0, 9.2, 0.38, "Survey Report / Field Research", 17, MID)
    text(slide, 0.82, 4.45, 5.8, 0.36, "汇报人：姓名", 15, DARK)
    text(slide, 0.82, 4.86, 5.8, 0.36, "单位：北京师范大学 ××学院", 15, DARK)
    text(slide, 0.82, 5.27, 5.8, 0.36, "调研时间：20XX年X月", 15, DARK)
    text(slide, 0.82, 6.95, 5.8, 0.22, "Research Briefing Template", 9, MID)


def agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "目录", 2)
    items = ["调研背景与目标", "调研对象与方法", "样本概况", "核心发现", "问题诊断与建议"]
    for i, item in enumerate(items, 1):
        y = 1.68 + (i - 1) * 0.78
        text(slide, 1.05, y, 0.45, 0.32, f"{i:02d}", 14, BLUE, True)
        text(slide, 1.75, y - 0.02, 7.5, 0.38, item, 20, DARK)


def section(prs, num, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    slide.shapes.add_picture(str(LOGO), Inches(11.65), Inches(0.36), width=Inches(0.84))
    rule(slide, 1.08)
    rule(slide, 6.86)
    text(slide, 0.86, 2.25, 1.1, 0.48, f"{num:02d}", 22, BLUE, True)
    text(slide, 2.05, 2.2, 8.2, 0.62, title, 34, BLUE, True)
    text(slide, 2.08, 2.95, 8.2, 0.35, subtitle, 16, MID)


def background(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "调研背景与目标", 4)
    text(slide, 0.95, 1.68, 4.7, 0.35, "背景", 18, BLUE, True)
    bullets(slide, 0.98, 2.18, 4.8, 3.6, [
        "说明调研缘起、现实情境与问题来源",
        "概括相关政策、项目或实践背景",
        "指出需要通过调研回答的关键问题"
    ])
    text(slide, 7.05, 1.68, 4.7, 0.35, "目标", 18, BLUE, True)
    bullets(slide, 7.08, 2.18, 4.8, 3.6, [
        "识别现状特征与主要差异",
        "发现痛点、需求和影响因素",
        "形成可执行的改进建议"
    ])


def method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "调研对象与方法", 5)
    heads = ["对象", "方法", "过程"]
    body = [
        "地区/学校/机构/人群；样本规模 N = ×××",
        "问卷、访谈、座谈、观察、二手资料分析",
        "方案设计、试测修订、正式调研、数据清洗"
    ]
    for i, head in enumerate(heads):
        x = 0.95 + i * 4.05
        text(slide, x, 1.98, 3.0, 0.34, head, 17, BLUE, True)
        text(slide, x, 2.58, 3.25, 1.5, body[i], 15, DARK)


def sample(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "样本概况", 6)
    chart_data = CategoryChartData()
    chart_data.categories = ["类别A", "类别B", "类别C", "类别D"]
    chart_data.add_series("占比", (32, 28, 24, 16))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(0.95), Inches(1.68), Inches(5.2), Inches(4.65),
        chart_data
    ).chart
    chart.has_legend = True
    text(slide, 7.0, 1.78, 3.8, 0.35, "样本说明", 18, BLUE, True)
    bullets(slide, 7.05, 2.32, 4.2, 3.15, [
        "样本覆盖范围与结构",
        "关键分组变量",
        "数据质量与有效样本说明"
    ], 15)


def findings(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "核心发现", 7)
    chart_data = CategoryChartData()
    chart_data.categories = ["指标1", "指标2", "指标3", "指标4"]
    chart_data.add_series("得分", (3.8, 4.2, 3.4, 4.5))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.95), Inches(1.68), Inches(6.4), Inches(4.65),
        chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = BLUE
    text(slide, 8.05, 1.78, 3.2, 0.34, "摘要", 18, BLUE, True)
    bullets(slide, 8.1, 2.32, 3.5, 2.9, [
        "发现一：总体趋势或主要特征",
        "发现二：不同群体存在差异",
        "发现三：关键需求较为集中"
    ], 15)


def interview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "访谈与现场材料", 8)
    pale_box(slide, 0.95, 1.68, 7.0, 4.55)
    text(slide, 0.95, 3.86, 7.0, 0.3, "现场照片 / 访谈摘录 / 材料截图占位", 16, MID, align=PP_ALIGN.CENTER)
    text(slide, 8.55, 1.78, 3.2, 0.34, "说明", 18, BLUE, True)
    bullets(slide, 8.6, 2.32, 3.5, 2.75, [
        "材料来源与调研时间",
        "与核心发现的关系",
        "典型案例或代表性观点"
    ], 15)


def problems(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "问题诊断", 9)
    items = [
        ("问题一", "描述现象、影响范围和可能原因。"),
        ("问题二", "说明群体差异、资源约束或流程堵点。"),
        ("问题三", "指出需要优先处理的风险或短板。"),
    ]
    for i, (head, body) in enumerate(items):
        y = 1.88 + i * 1.18
        text(slide, 1.0, y, 2.1, 0.32, head, 17, BLUE, True)
        text(slide, 3.05, y, 7.7, 0.42, body, 15, DARK)


def recommendations(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "对策建议", 10)
    items = [
        ("短期行动", "明确 1-3 项可立即推进的措施。"),
        ("中期优化", "完善机制、流程、培训或资源配置。"),
        ("长期建设", "形成制度化、可评估、可持续的改进方案。"),
        ("评估指标", "列出后续跟踪的数据指标与责任主体。"),
    ]
    for i, (head, body) in enumerate(items):
        y = 1.72 + i * 1.0
        text(slide, 1.0, y, 2.1, 0.32, head, 16, BLUE, True)
        text(slide, 3.05, y, 7.7, 0.36, body, 15, DARK)


def thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    slide.shapes.add_picture(str(LOGO), Inches(11.62), Inches(0.36), width=Inches(0.86))
    rule(slide, 1.08)
    rule(slide, 6.86)
    text(slide, 0.88, 2.55, 7.6, 0.72, "谢谢聆听", 40, BLUE, True)
    text(slide, 0.9, 3.42, 7.6, 0.38, "敬请批评指正", 20, MID)
    text(slide, 0.92, 5.25, 5.8, 0.32, "姓名 | 北京师范大学 ××学院", 14, DARK)
    text(slide, 0.92, 5.65, 5.8, 0.32, "email@bnu.edu.cn", 14, DARK)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    agenda(prs)
    section(prs, 1, "调研背景", "Research Background and Objectives")
    background(prs)
    method(prs)
    sample(prs)
    findings(prs)
    interview(prs)
    problems(prs)
    recommendations(prs)
    thanks(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
