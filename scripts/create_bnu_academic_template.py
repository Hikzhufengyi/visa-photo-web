from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT = Path("/Users/seafon/Downloads/北京师范大学_学术汇报模板.pptx")
LOGO = Path("/Users/seafon/Downloads/bnu_logo_transparent.png")

BLUE = RGBColor(0, 91, 172)
DARK = RGBColor(34, 43, 54)
MID = RGBColor(100, 110, 122)
PALE = RGBColor(248, 250, 252)
SOFT = RGBColor(234, 241, 249)
WHITE = RGBColor(255, 255, 255)


def text(slide, x, y, w, h, value, size=18, color=DARK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def bullets(slide, x, y, w, h, values, size=17, color=DARK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, value in enumerate(values):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = value
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(12)
    return box


def logo(slide, width=3.0):
    slide.shapes.add_picture(str(LOGO), Inches(13.333 - width - 0.55), Inches(0.28), width=Inches(width))


def header(slide, title, page=None):
    logo(slide, 2.85)
    text(slide, 0.68, 0.46, 7.6, 0.42, title, 22, BLUE, True)
    if page:
        text(slide, 12.15, 7.05, 0.45, 0.2, str(page), 8.5, MID, align=PP_ALIGN.RIGHT)


def pale_box(slide, x, y, w, h):
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = PALE
    box.line.fill.background()
    return box


def cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    logo(slide, 3.9)
    text(slide, 0.78, 2.15, 9.2, 0.7, "学术汇报题目", 36, BLUE, True)
    text(slide, 0.8, 2.95, 9.2, 0.38, "Subtitle / Research Topic", 18, MID)
    text(slide, 0.82, 4.45, 5.3, 0.36, "汇报人：姓名", 15, DARK)
    text(slide, 0.82, 4.86, 5.3, 0.36, "单位：北京师范大学 ××学院", 15, DARK)
    text(slide, 0.82, 5.27, 5.3, 0.36, "日期：2026年", 15, DARK)
    text(slide, 0.82, 6.95, 5.2, 0.22, "Academic Presentation", 9, MID)


def agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "目录", 2)
    items = ["研究背景与问题提出", "文献综述与理论框架", "研究设计与方法", "实证结果与分析", "结论与展望"]
    for i, item in enumerate(items, 1):
        y = 1.55 + (i - 1) * 0.82
        text(slide, 1.05, y, 0.45, 0.32, f"{i:02d}", 14, BLUE, True)
        text(slide, 1.75, y - 0.02, 7.5, 0.38, item, 20, DARK)


def section(prs, num, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PALE
    logo(slide, 3.25)
    text(slide, 0.86, 2.25, 1.1, 0.48, f"{num:02d}", 22, BLUE, True)
    text(slide, 2.05, 2.2, 8.2, 0.62, title, 34, BLUE, True)
    text(slide, 2.08, 2.95, 8.2, 0.35, subtitle, 16, MID)


def background_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "研究背景与问题提出", 4)
    text(slide, 0.9, 1.55, 4.7, 0.35, "问题背景", 18, BLUE, True)
    bullets(slide, 0.95, 2.05, 4.85, 3.8, [
        "概述研究对象、现实背景与实践情境",
        "说明已有研究或实践中的关键缺口",
        "明确本研究关注的核心问题"
    ], 16)
    text(slide, 7.0, 1.55, 4.7, 0.35, "研究问题", 18, BLUE, True)
    bullets(slide, 7.05, 2.05, 4.75, 3.8, [
        "RQ1：核心问题一",
        "RQ2：核心问题二",
        "RQ3：机制与边界条件"
    ], 16)


def framework(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "理论框架", 5)
    positions = [("前因变量", 1.05, 2.4), ("中介机制", 4.25, 2.4), ("结果变量", 7.45, 2.4)]
    for label, x, y in positions:
        pale_box(slide, x, y, 2.3, 0.82)
        text(slide, x, y + 0.24, 2.3, 0.25, label, 16, BLUE, True, PP_ALIGN.CENTER)
    for x1, x2 in [(3.35, 4.18), (6.55, 7.38)]:
        line = slide.shapes.add_connector(1, Inches(x1), Inches(2.81), Inches(x2), Inches(2.81))
        line.line.color.rgb = BLUE
        line.line.width = Pt(1.2)
    text(slide, 4.25, 4.35, 2.3, 0.32, "调节变量", 16, BLUE, True, PP_ALIGN.CENTER)
    bullets(slide, 0.95, 5.55, 10.6, 0.8, ["用一句话说明理论依据、变量关系和可替代解释。"], 15, MID)


def method(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "研究设计与方法", 6)
    heads = ["样本与数据", "变量测量", "模型设定"]
    body = [
        "样本来源、样本规模、时间范围",
        "核心变量、控制变量、信度效度",
        "识别策略、估计方法、稳健性检验"
    ]
    for i, head in enumerate(heads):
        x = 0.95 + i * 4.05
        text(slide, x, 1.85, 3.0, 0.34, head, 17, BLUE, True)
        text(slide, x, 2.45, 3.25, 1.4, body[i], 15, DARK)


def chart_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "结果展示", 7)
    chart_data = CategoryChartData()
    chart_data.categories = ["模型1", "模型2", "模型3", "模型4"]
    chart_data.add_series("估计值", (0.18, 0.31, 0.27, 0.43))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.95), Inches(1.55), Inches(6.4), Inches(4.8),
        chart_data
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = False
    chart.plots[0].series[0].format.fill.solid()
    chart.plots[0].series[0].format.fill.fore_color.rgb = BLUE
    text(slide, 8.05, 1.65, 3.2, 0.34, "主要发现", 18, BLUE, True)
    bullets(slide, 8.1, 2.18, 3.5, 2.9, [
        "核心变量方向与理论预期一致",
        "机制检验支持主要假设",
        "稳健性检验结论稳定"
    ], 15)


def table_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "表格页", 8)
    rows, cols = 6, 5
    table = slide.shapes.add_table(rows, cols, Inches(0.9), Inches(1.65), Inches(11.4), Inches(4.55)).table
    for c, name in enumerate(["变量", "模型1", "模型2", "模型3", "模型4"]):
        cell = table.cell(0, c)
        cell.text = name
        cell.fill.solid()
        cell.fill.fore_color.rgb = SOFT
        for p in cell.text_frame.paragraphs:
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = BLUE
            p.alignment = PP_ALIGN.CENTER
    data = [
        ["核心变量 X", "0.183**", "0.205***", "0.196***", "0.211***"],
        ["中介变量 M", "", "0.142**", "0.137**", "0.129**"],
        ["调节项 X×W", "", "", "0.084*", "0.091*"],
        ["控制变量", "是", "是", "是", "是"],
        ["观测值", "1,280", "1,280", "1,280", "1,280"],
    ]
    for r, row in enumerate(data, 1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(11)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
    text(slide, 0.92, 6.45, 9.5, 0.28, "注：* p<0.10, ** p<0.05, *** p<0.01。", 10, MID)


def image_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "图片与案例材料", 9)
    pale_box(slide, 0.95, 1.55, 7.0, 4.75)
    text(slide, 0.95, 3.8, 7.0, 0.3, "图片 / 案例截图占位", 16, MID, align=PP_ALIGN.CENTER)
    text(slide, 8.55, 1.65, 3.2, 0.34, "说明", 18, BLUE, True)
    bullets(slide, 8.6, 2.18, 3.5, 2.75, [
        "材料来源与时间",
        "与研究问题的关联",
        "关键观察或解释"
    ], 15)


def conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    header(slide, "结论与展望", 10)
    items = [
        ("主要结论", "概括最重要的 2-3 个发现。"),
        ("理论贡献", "说明对既有理论或机制解释的增量。"),
        ("实践启示", "面向政策、教育实践或管理给出建议。"),
        ("局限与展望", "说明数据、方法和外部效度边界。"),
    ]
    for i, (head, body) in enumerate(items):
        y = 1.6 + i * 1.05
        text(slide, 1.0, y, 2.1, 0.32, head, 16, BLUE, True)
        text(slide, 3.05, y, 7.6, 0.36, body, 15, DARK)


def thanks(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE
    logo(slide, 3.9)
    text(slide, 0.88, 2.55, 7.6, 0.72, "谢谢聆听", 40, BLUE, True)
    text(slide, 0.9, 3.42, 7.6, 0.38, "敬请批评指正", 20, MID)
    text(slide, 0.92, 5.25, 5.8, 0.32, "姓名 | 北京师范大学 ××学院", 14, DARK)
    text(slide, 0.92, 5.65, 5.8, 0.32, "email@bnu.edu.cn", 14, DARK)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    cover(prs)
    agenda(prs)
    section(prs, 1, "研究背景", "Research Background and Questions")
    background_slide(prs)
    framework(prs)
    method(prs)
    chart_slide(prs)
    table_slide(prs)
    image_slide(prs)
    conclusion(prs)
    thanks(prs)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    main()
